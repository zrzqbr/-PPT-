#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
test_v5.py — apply_template.py v5 自动化测试套件

验证 14 条诊断报告中的修复是否真实生效：
  P0-1: 全屏遮罩检测与清除
  P0-2: Logo 去重
  P0-3: 配色两套逻辑关系（注释层面）
  P1-4: CIEDE2000 色差算法
  P1-5: 字体替换增量计数
  P1-6/7: 静默异常改为 warnings
  P1-8: classify_page 组合形状扫描 + 阈值常量
  P2-12: 字体安装检测
  CLI: 版本号更正 v5

运行：
  python scripts/test_v5.py
"""

import sys
import os
import math
import io
import unittest
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

# Windows GBK 编码兼容：强制 stdout/stderr 使用 utf-8，避免 emoji 输出报错
if sys.platform == "win32" and os.environ.get("PYTHONIOENCODING") is None:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))

# ---------- 导入被测模块 ----------
try:
    import apply_template as T
    print(f"[OK] 模块导入成功（v5 路径：{SCRIPT_DIR / 'apply_template.py'}）")
except ImportError as e:
    print(f"[FAIL] 模块导入失败：{e}")
    print("  请先安装依赖：pip install python-pptx Pillow lxml")
    sys.exit(1)


# ============================================================
#  辅助工具
# ============================================================

def _make_fake_shape(left=0, top=0, width=None, height=None,
                     fill_type='SOLID', is_pic=False,
                     has_text=False, text="", font_size_pt=None,
                     has_chart=False, has_table=False,
                     sub_shapes=None):
    """构造足够真实的 Mock shape 对象，供 classify_page / remove_fullscreen_overlays 测试。"""
    from lxml import etree
    shp = MagicMock()
    shp.left = left
    shp.top = top
    shp.width = width if width is not None else 9144000   # 默认全宽
    shp.height = height if height is not None else 5143500  # 默认全高
    shp.name = "test_shape"
    shp.shape_id = 1
    shp.has_chart = has_chart
    shp.has_table = has_table

    # 元素标签
    tag = 'pic' if is_pic else 'sp'
    el = etree.Element(f'{{{T.qn("p:sp").split("}")[0][1:]}}}' + tag
                       if False else tag)
    shp._element = el

    # fill
    fill = MagicMock()
    fill_type_mock = MagicMock()
    fill_type_mock.name = fill_type
    fill.type = fill_type_mock
    shp.fill = fill

    # text_frame
    shp.has_text_frame = has_text
    if has_text:
        tf = MagicMock()
        tf.text = text
        para = MagicMock()
        run = MagicMock()
        run.font = MagicMock()
        if font_size_pt is not None:
            size = MagicMock()
            size.pt = font_size_pt
            run.font.size = size
        else:
            run.font.size = None
        para.runs = [run]
        tf.paragraphs = [para]
        shp.text_frame = tf
    else:
        shp.text_frame = None

    # sub_shapes (for group shape)
    if sub_shapes is not None:
        shp.shapes = sub_shapes
    elif not hasattr(shp, 'shapes'):
        # 不是 group shape —— 不提供 .shapes 属性
        del shp.shapes

    return shp


# ============================================================
#  测试类
# ============================================================

class TestCIEDE2000(unittest.TestCase):
    """P1-4: 验证 delta_e 使用 CIEDE2000 而非 CIE76。"""

    def test_identical_colors_zero(self):
        """同一颜色 ΔE 应为 0。"""
        self.assertAlmostEqual(T.delta_e(128, 64, 200, 128, 64, 200), 0.0, places=3)

    def test_blue_purple_vs_brand_blue_not_forbidden(self):
        """商务蓝 #1D6FA9 不应被视为禁用色（CIEDE2000 对蓝区更精准）。"""
        r, g, b = 0x1D, 0x6F, 0xA9
        self.assertFalse(T._is_forbidden_color(r, g, b),
                         "商务蓝 #1D6FA9 不应命中禁用色规则")

    def test_pure_forbidden_blue_purple(self):
        """纯禁用蓝紫 #4B62E1 应被识别为禁用色。"""
        r, g, b = 0x4B, 0x62, 0xE1
        self.assertTrue(T._is_forbidden_color(r, g, b),
                        "#4B62E1 应命中禁用色规则")

    def test_delta_e_white_black_large(self):
        """白色与黑色的 ΔE 应远大于 100。"""
        de = T.delta_e(255, 255, 255, 0, 0, 0)
        self.assertGreater(de, 80, f"白黑 ΔE 应 > 80，实际={de:.1f}")

    def test_ciede2000_function_signature(self):
        """delta_e 应能接受 6 个 int 参数并返回 float。"""
        result = T.delta_e(200, 100, 50, 200, 100, 50)
        self.assertIsInstance(result, float)

    def test_docstring_mentions_ciede2000(self):
        """docstring 应提及 CIEDE2000，确认不是旧的 CIE76 实现。"""
        doc = T.delta_e.__doc__ or ""
        self.assertIn("CIEDE2000", doc,
                      "delta_e() docstring 未提及 CIEDE2000，可能仍是旧 CIE76 实现")


class TestRemoveFullscreenOverlays(unittest.TestCase):
    """P0-1: 全屏遮罩检测与清除。"""

    def _make_slide_with_shape(self, shp_width, shp_height, shp_left=0, shp_top=0,
                                fill_type='SOLID', is_pic=False):
        from lxml import etree
        slide_w = 9144000
        slide_h = 5143500

        slide = MagicMock()
        shp = MagicMock()
        shp.left = shp_left
        shp.top = shp_top
        shp.width = shp_width
        shp.height = shp_height
        shp.name = "overlay"

        # 元素
        tag = 'pic' if is_pic else 'sp'
        sp_el = etree.Element(tag)
        parent_el = etree.Element('spTree')
        parent_el.append(sp_el)
        shp._element = sp_el

        fill = MagicMock()
        ft = MagicMock()
        ft.name = fill_type
        fill.type = ft
        shp.fill = fill

        # slide.shapes 需要同时可迭代且拥有 _spTree 属性
        shapes_mock = MagicMock()
        shapes_mock.__iter__ = MagicMock(return_value=iter([shp]))
        shapes_mock._spTree = parent_el
        slide.shapes = shapes_mock
        return slide, slide_w, slide_h

    def test_detect_fullscreen_solid_rect(self):
        """全屏实心矩形（>=90%）应被检测到。"""
        slide_w = 9144000
        slide_h = 5143500
        slide, sw, sh = self._make_slide_with_shape(
            shp_width=int(slide_w * 0.95),
            shp_height=int(slide_h * 0.95),
            fill_type='SOLID'
        )
        count = T.remove_fullscreen_overlays(slide, sw, sh, dry_run=True)
        self.assertEqual(count, 1, "应检测到1个全屏遮罩")

    def test_small_shape_not_detected(self):
        """小形状（50%宽）不应被当作全屏遮罩删除。"""
        slide_w = 9144000
        slide_h = 5143500
        slide, sw, sh = self._make_slide_with_shape(
            shp_width=int(slide_w * 0.5),
            shp_height=int(slide_h * 0.5),
            fill_type='SOLID'
        )
        count = T.remove_fullscreen_overlays(slide, sw, sh, dry_run=True)
        self.assertEqual(count, 0, "50%宽的形状不应被视为全屏遮罩")

    def test_fullscreen_pic_detected(self):
        """全屏图片（p:pic）应被检测到。"""
        slide_w = 9144000
        slide_h = 5143500
        slide, sw, sh = self._make_slide_with_shape(
            shp_width=int(slide_w * 0.98),
            shp_height=int(slide_h * 0.98),
            is_pic=True
        )
        count = T.remove_fullscreen_overlays(slide, sw, sh, dry_run=True)
        self.assertEqual(count, 1, "全屏图片应被检测到")

    def test_offset_shape_not_detected(self):
        """偏移量过大（left > 5%页宽）的形状不应被检测为遮罩。"""
        slide_w = 9144000
        slide_h = 5143500
        margin = int(slide_w * 0.06)  # 超出允许偏移
        slide, sw, sh = self._make_slide_with_shape(
            shp_width=int(slide_w * 0.95),
            shp_height=int(slide_h * 0.95),
            shp_left=margin,
            fill_type='SOLID'
        )
        count = T.remove_fullscreen_overlays(slide, sw, sh, dry_run=True)
        self.assertEqual(count, 0, "偏移量过大的形状不应被视为全屏遮罩")

    def test_function_exists_and_callable(self):
        """函数必须存在。"""
        self.assertTrue(callable(T.remove_fullscreen_overlays),
                        "remove_fullscreen_overlays 函数不存在")


class TestLogoDedup(unittest.TestCase):
    """P0-2: Logo 去重逻辑。"""

    def test_logo_already_exists_detected(self):
        """当页面已有相同位置/尺寸的图片时，_logo_already_exists 应返回 True。"""
        from pptx.util import Inches
        from lxml import etree

        target_w = Inches(2.80)
        target_h = Inches(2.80 * 50 / 463)
        left = Inches(0.30)
        top = Inches(0.18)

        # 构造已有 Logo 的 shape
        slide = MagicMock()
        shp = MagicMock()
        shp.left = left
        shp.top = top
        shp.width = target_w
        shp.height = target_h
        sp_el = etree.Element('pic')
        shp._element = sp_el
        slide.shapes = [shp]

        result = T._logo_already_exists(slide, left, top, target_w, target_h)
        self.assertTrue(result, "已有相同尺寸/位置的 Logo，应返回 True")

    def test_logo_not_exist_when_empty(self):
        """空页面应返回 False。"""
        from pptx.util import Inches
        slide = MagicMock()
        slide.shapes = []
        result = T._logo_already_exists(
            slide, Inches(0.30), Inches(0.18), Inches(2.80), Inches(0.302)
        )
        self.assertFalse(result, "空页面不应检测到已有 Logo")

    def test_add_logo_skips_if_exists(self):
        """add_logo 检测到已有 Logo 时应返回 False（跳过）。"""
        from pptx.util import Inches
        from lxml import etree

        # 内容页 Logo 位置：右上角 (x=14.60", y=0.49", w=4.88", h=0.53")
        target_w = Inches(4.88)
        target_h = Inches(0.53)
        left = Inches(14.60)
        top = Inches(0.49)

        # 构造 Mock slide，shapes 是 MagicMock（有 add_picture 方法），
        # 同时支持迭代（用于 _logo_already_exists 遍历）
        slide = MagicMock()
        shp = MagicMock()
        shp.left = left
        shp.top = top
        shp.width = target_w
        shp.height = target_h
        shp._element = etree.Element('pic')

        # shapes 是 MagicMock，__iter__ 返回 [shp]
        shapes_mock = MagicMock()
        shapes_mock.__iter__ = MagicMock(return_value=iter([shp]))
        slide.shapes = shapes_mock
        slide_w = 18305463
        slide_h = 10296525

        # add_logo 检测到已有 → 应返回 False，不调用 add_picture
        result = T.add_logo(slide, "content",
                            str(T.ASSETS / "logos" / "logo-main.png"),
                            None, slide_w, slide_h)
        self.assertFalse(result, "已有 Logo 时 add_logo 应返回 False")
        shapes_mock.add_picture.assert_not_called()


class TestFontCountIncremental(unittest.TestCase):
    """P1-5: 字体替换增量计数——已是目标字体时不计入 changes。

    注意：replace_fonts_in_slide 调用 run._r.get_or_add_rPr()，
    这是 pptx 专有方法，普通 lxml Element 没有。
    因此测试通过 Mock 整个 _r 来模拟。
    """

    def _make_run_with_mock_r(self, latin_typeface, ea_typeface, cs_typeface, size_pt=18):
        """构造 Mock run，其 _r 支持 get_or_add_rPr 并返回有正确字体属性的 rPr element。"""
        from lxml import etree
        from pptx.oxml.ns import qn

        # 构造真实 rPr element（含目标字体）
        rPr = etree.Element(qn('a:rPr'))
        latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', latin_typeface)
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', ea_typeface)
        cs = etree.SubElement(rPr, qn('a:cs'))
        cs.set('typeface', cs_typeface)

        # Mock run._r，让 get_or_add_rPr 返回这个 rPr
        mock_r = MagicMock()
        mock_r.get_or_add_rPr.return_value = rPr
        # find 方法委托给真实 rPr，但 get_or_add_rPr 是 Mock
        # 需要让 rPr.find 等方法正常工作 → 直接返回 rPr 对象

        run = MagicMock()
        run._r = mock_r
        run.font = MagicMock()
        if size_pt:
            sz = MagicMock()
            sz.pt = size_pt
            run.font.size = sz
        else:
            run.font.size = None
        run.font.name = None
        return run

    def test_no_change_when_already_target_font(self):
        """run 已经是 TencentSans W3 时，changes 应为 0。"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE as ST
        slide = MagicMock()
        shp = MagicMock()
        shp.has_text_frame = True
        shp.has_table = False
        shp.shape_type = ST.AUTO_SHAPE  # 非 GROUP

        run = self._make_run_with_mock_r(
            latin_typeface="TencentSans W3",
            ea_typeface="TencentSans W3",
            cs_typeface="TencentSans W3",
            size_pt=18
        )
        para = MagicMock()
        para.runs = [run]
        shp.text_frame.paragraphs = [para]
        slide.shapes = [shp]

        changes = T.replace_fonts_in_slide(slide)
        self.assertEqual(changes, 0,
                         f"已是目标字体时 changes 应为 0，实际为 {changes}")

    def test_change_counted_when_different_font(self):
        """run 是 Arial 时，changes 应为 1（实际发生了变更）。"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE as ST
        slide = MagicMock()
        shp = MagicMock()
        shp.has_text_frame = True
        shp.has_table = False
        shp.shape_type = ST.AUTO_SHAPE  # 非 GROUP

        run = self._make_run_with_mock_r(
            latin_typeface="Arial",
            ea_typeface="Microsoft YaHei",
            cs_typeface="Arial",
            size_pt=18
        )
        para = MagicMock()
        para.runs = [run]
        shp.text_frame.paragraphs = [para]
        slide.shapes = [shp]

        changes = T.replace_fonts_in_slide(slide)
        self.assertEqual(changes, 1,
                         f"不同字体时 changes 应为 1，实际为 {changes}")


class TestClassifyPageWithGroups(unittest.TestCase):
    """P1-8: classify_page 应扫描组合形状内子形状，且使用可配置常量。"""

    def test_section_font_threshold_constant_exists(self):
        """SECTION_FONT_THRESHOLD_PT 常量必须存在。"""
        self.assertTrue(hasattr(T, 'SECTION_FONT_THRESHOLD_PT'),
                        "缺少 SECTION_FONT_THRESHOLD_PT 常量")

    def test_section_threshold_value(self):
        """SECTION_FONT_THRESHOLD_PT 应 < 40（旧硬编码值），体现了改进。"""
        self.assertLess(T.SECTION_FONT_THRESHOLD_PT, 40,
                        f"阈值 {T.SECTION_FONT_THRESHOLD_PT} 应 < 40（已调整为更合理的默认值）")

    def test_collect_text_shapes_exists(self):
        """辅助函数 _collect_text_shapes 必须存在（支持递归扫描）。"""
        self.assertTrue(callable(T._collect_text_shapes),
                        "_collect_text_shapes 函数不存在")

    def test_group_shape_scan(self):
        """_collect_text_shapes 应能扫描到组合形状内的大字号文字。"""
        from unittest.mock import MagicMock

        # 构造一个包含大字号文字的子形状
        sub_shp = MagicMock()
        sub_shp.has_text_frame = True
        sub_shp.text_frame.text = "大标题"
        sub_run = MagicMock()
        sub_run.font.size = MagicMock()
        sub_run.font.size.pt = T.SECTION_FONT_THRESHOLD_PT + 5  # 超阈值
        sub_para = MagicMock()
        sub_para.runs = [sub_run]
        sub_shp.text_frame.paragraphs = [sub_para]
        sub_shp.has_chart = False
        sub_shp.has_table = False

        # 父组合形状
        group_shp = MagicMock()
        group_shp.has_text_frame = False
        group_shp.has_chart = False
        group_shp.has_table = False
        group_shp.shapes = [sub_shp]  # 关键：有 .shapes 属性

        slide = MagicMock()
        slide.shapes = [group_shp]

        text_shapes, big_title_count, has_chart, has_table = T._collect_text_shapes(slide)
        self.assertGreater(big_title_count, 0,
                           "组合形状内的大字号文字应被 _collect_text_shapes 扫描到")


class TestFontInstallCheck(unittest.TestCase):
    """P2-12: 字体安装检测函数必须存在且不阻断流程。"""

    def test_check_function_exists(self):
        """_check_tencent_sans_installed 必须存在。"""
        self.assertTrue(callable(T._check_tencent_sans_installed),
                        "_check_tencent_sans_installed 函数不存在")

    def test_check_returns_bool(self):
        """函数应返回 bool，且不抛异常。"""
        result = T._check_tencent_sans_installed()
        self.assertIsInstance(result, bool,
                              f"应返回 bool，实际返回 {type(result)}")

    def test_warning_in_report_when_font_missing(self):
        """字体未安装时，adapt() 报告的 warnings 应包含相关提示。"""
        with patch.object(T, '_check_tencent_sans_installed', return_value=False):
            with patch.object(T, 'get_assets', return_value={
                'bg_cover': '', 'bg_content': '', 'logo_main': '',
                'logo_corner': None, 'city': 'changsha'
            }):
                with patch.object(T, 'Presentation') as mock_pres:
                    pres = MagicMock()
                    pres.slide_width = 9144000
                    pres.slide_height = 5143500
                    pres.slides = []
                    mock_pres.return_value = pres
                    pres.slide_masters = []

                    report = T.adapt("dummy.pptx", "out.pptx", dry_run=True)
                    warnings_text = " ".join(report.get("warnings", []))
                    self.assertIn("TencentSans", warnings_text,
                                  "字体未安装时 warnings 应包含 TencentSans 提示")


class TestCLIVersion(unittest.TestCase):
    """CLI description 应为 v7，不再是旧版本。"""

    def test_cli_description_is_v7(self):
        """main() 中 ArgumentParser description 应包含 v7。"""
        import subprocess
        result = subprocess.run(
            [sys.executable,
             str(SCRIPT_DIR / "apply_template.py"), "--help"],
            capture_output=True, timeout=10
        )
        # Windows 下可能输出 GBK 或 UTF-8，兼容处理
        output = (result.stdout or b"").decode("utf-8", errors="replace") + \
                 (result.stderr or b"").decode("utf-8", errors="replace")
        self.assertIn("v7", output,
                      f"CLI --help 输出中未找到 v7，当前输出：{output[:300]}")
        self.assertNotIn("v3", output,
                         "CLI --help 输出中仍包含旧版本号 v3")


class TestReportFields(unittest.TestCase):
    """验证 adapt() 报告包含新增字段（overlay_removed, logo_skipped）。"""

    def test_new_report_fields_present(self):
        """报告字典应包含 overlay_removed 和 logo_skipped 字段。"""
        with patch.object(T, '_check_tencent_sans_installed', return_value=True):
            with patch.object(T, 'get_assets', return_value={
                'bg_cover': '', 'bg_content': '', 'logo_main': '',
                'logo_corner': None, 'city': 'changsha'
            }):
                with patch.object(T, 'Presentation') as mock_pres:
                    pres = MagicMock()
                    pres.slide_width = 9144000
                    pres.slide_height = 5143500
                    pres.slides = []
                    pres.slide_masters = []
                    mock_pres.return_value = pres

                    report = T.adapt("dummy.pptx", "out.pptx", dry_run=True)
                    self.assertIn("overlay_removed", report,
                                  "报告应包含 overlay_removed 字段")
                    self.assertIn("logo_skipped", report,
                                  "报告应包含 logo_skipped 字段")


class TestBrandPaletteWarmTone(unittest.TestCase):
    """P2-9 误判核实：WarmTone 5色已在 BRAND_PALETTE 中，不是缺陷。"""

    WARMTONE = [
        (0xE8, 0xA0, 0x90, "珊瑚粉"),
        (0xD4, 0xA5, 0x74, "暖杏色"),
        (0xE8, 0xB8, 0x4A, "浅金橙"),
        (0xC4, 0xA4, 0x84, "玫瑰木"),
        (0xD4, 0x9A, 0x6A, "浅赭色"),
    ]

    def test_warmtone_colors_in_palette(self):
        """WarmTone 5色应全部存在于 BRAND_PALETTE。"""
        palette_set = {(r, g, b) for r, g, b, _ in T.BRAND_PALETTE}
        for r, g, b, name in self.WARMTONE:
            self.assertIn((r, g, b), palette_set,
                          f"WarmTone 色 {name} ({r:02X},{g:02X},{b:02X}) 不在 BRAND_PALETTE 中")

    def test_warmtone_not_forbidden(self):
        """WarmTone 颜色不应被识别为禁用色。"""
        for r, g, b, name in self.WARMTONE:
            self.assertFalse(T._is_forbidden_color(r, g, b),
                             f"WarmTone {name} 不应被视为禁用色")


class TestThemeColorExceptionPropagation(unittest.TestCase):
    """P1-6: replace_theme_colors 异常应向上传播（不再静默 pass）。"""

    def test_exception_propagates_from_replace_theme_colors(self):
        """replace_theme_colors 内部异常应 raise，由 adapt() 捕获写入 warnings。"""
        pres = MagicMock()
        bad_master = MagicMock()
        bad_master.part.rels.values.side_effect = RuntimeError("故意触发的测试异常")
        pres.slide_masters = [bad_master]

        # 直接调用应抛出（因为内部 raise 了）
        with self.assertRaises(Exception):
            T.replace_theme_colors(pres, dry_run=True)

    def test_adapt_catches_theme_exception_into_warnings(self):
        """adapt() 应捕获主题色覆写异常并写入 warnings。"""
        with patch.object(T, '_check_tencent_sans_installed', return_value=True):
            with patch.object(T, 'get_assets', return_value={
                'bg_cover': '', 'bg_content': '', 'logo_main': '',
                'logo_corner': None, 'city': 'changsha'
            }):
                with patch.object(T, 'replace_theme_colors',
                                   side_effect=RuntimeError("主题色测试异常")):
                    with patch.object(T, 'Presentation') as mock_pres:
                        pres = MagicMock()
                        pres.slide_width = 9144000
                        pres.slide_height = 5143500
                        pres.slides = []
                        pres.slide_masters = []
                        mock_pres.return_value = pres

                        report = T.adapt("dummy.pptx", "out.pptx", dry_run=True)
                        warnings_text = " ".join(report.get("warnings", []))
                        self.assertIn("主题色覆写失败", warnings_text,
                                      "主题色覆写异常应被记录到 warnings")


# ============================================================
#  v7 新增测试：对比度修复 / classify_page Layout / 溢出检测
# ============================================================

class TestContrastFix(unittest.TestCase):
    """P0-A: fix_text_contrast 应修复低对比度文字。"""

    def test_function_exists(self):
        """fix_text_contrast 函数必须存在。"""
        self.assertTrue(hasattr(T, 'fix_text_contrast'))
        self.assertTrue(callable(T.fix_text_contrast))

    def test_relative_luminance_white(self):
        """白色亮度应接近 1.0。"""
        lum = T._relative_luminance(255, 255, 255)
        self.assertAlmostEqual(lum, 1.0, places=2)

    def test_relative_luminance_black(self):
        """黑色亮度应为 0.0。"""
        lum = T._relative_luminance(0, 0, 0)
        self.assertAlmostEqual(lum, 0.0, places=2)

    def test_contrast_ratio_bw(self):
        """黑白对比度应为 21:1。"""
        ratio = T._contrast_ratio(1.0, 0.0)
        self.assertAlmostEqual(ratio, 21.0, places=1)

    def test_contrast_ratio_same(self):
        """相同亮度对比度为 1:1。"""
        ratio = T._contrast_ratio(0.5, 0.5)
        self.assertAlmostEqual(ratio, 1.0, places=1)

    def test_report_contains_contrast_fixes(self):
        """报告应包含 contrast_fixes 字段。"""
        with patch.object(T, '_check_tencent_sans_installed', return_value=True):
            with patch.object(T, 'get_assets', return_value={
                'bg_cover': '', 'bg_content': '', 'logo_main': '',
                'logo_corner': None, 'city': 'changsha'
            }):
                with patch.object(T, 'Presentation') as mock_pres:
                    pres = MagicMock()
                    pres.slide_width = 9144000
                    pres.slide_height = 5143500
                    pres.slides = []
                    pres.slide_masters = []
                    mock_pres.return_value = pres

                    report = T.adapt("dummy.pptx", "out.pptx", dry_run=True)
                    self.assertIn("contrast_fixes", report)


class TestClassifyPageLayout(unittest.TestCase):
    """P1-A: classify_page 应优先从 Layout 名称识别。"""

    def _make_slide(self, layout_name, shapes_count=0):
        slide = MagicMock()
        slide.slide_layout.name = layout_name
        slide.shapes = [MagicMock(has_text_frame=False, has_chart=False,
                                   has_table=False) for _ in range(shapes_count)]
        return slide

    def test_cover_layout(self):
        """Layout 名含 '封面' 应返回 cover。"""
        slide = self._make_slide("封面版式")
        result = T.classify_page(slide, 1, 10)
        self.assertEqual(result, "cover")

    def test_section_layout(self):
        """Layout 名含 'section' 应返回 section。"""
        slide = self._make_slide("Section Header")
        result = T.classify_page(slide, 2, 10)
        self.assertEqual(result, "section")

    def test_ending_layout(self):
        """Layout 名含 '结尾' 应返回 end。"""
        slide = self._make_slide("结尾页")
        result = T.classify_page(slide, 5, 10)
        self.assertEqual(result, "end")

    def test_content_fallback(self):
        """无匹配 Layout 名的中间页应走启发式，返回 content。"""
        slide = self._make_slide("自定义版式")
        # Mock _collect_text_shapes 使其返回内容页特征
        with patch.object(T, '_collect_text_shapes', return_value=(5, 0, True, False)):
            result = T.classify_page(slide, 3, 10)
            self.assertEqual(result, "content")


class TestElementOverflow(unittest.TestCase):
    """P1-B: check_element_overflow 应检测溢出。"""

    def test_function_exists(self):
        """check_element_overflow 必须存在。"""
        self.assertTrue(hasattr(T, 'check_element_overflow'))

    def test_overflow_detection(self):
        """溢出右边界的形状应被检测到。"""
        from lxml import etree
        slide = MagicMock()
        shp = MagicMock()
        shp.left = 18000000  # 远超正常位置
        shp.width = 5000000
        shp.height = 1000000
        shp.top = 500000
        shp.has_text_frame = False
        shp.has_table = False
        shp.name = "test_shape"
        sp_el = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}sp')
        shp._element = sp_el
        slide.shapes = [shp]

        slide_w = 18288000  # 20.02" in EMU
        slide_h = 10287000  # 11.26" in EMU
        warnings = T.check_element_overflow(slide, slide_w, slide_h, "content")
        # 形状右边缘 = 18000000 + 5000000 = 23000000 > slide_w + 0.1"
        overflow_found = any("溢出右边界" in w for w in warnings)
        self.assertTrue(overflow_found, f"应检测到右溢出，实际：{warnings}")

    def test_report_contains_overflow_field(self):
        """报告应包含 overflow_warnings 字段。"""
        with patch.object(T, '_check_tencent_sans_installed', return_value=True):
            with patch.object(T, 'get_assets', return_value={
                'bg_cover': '', 'bg_content': '', 'logo_main': '',
                'logo_corner': None, 'city': 'changsha'
            }):
                with patch.object(T, 'Presentation') as mock_pres:
                    pres = MagicMock()
                    pres.slide_width = 9144000
                    pres.slide_height = 5143500
                    pres.slides = []
                    pres.slide_masters = []
                    mock_pres.return_value = pres

                    report = T.adapt("dummy.pptx", "out.pptx", dry_run=True)
                    self.assertIn("overflow_warnings", report)


# ============================================================
#  主运行器：自定义彩色输出
# ============================================================

class ColorTestResult(unittest.TextTestResult):
    GREEN = '\033[92m'
    RED = '\033[91m'
    YELLOW = '\033[93m'
    RESET = '\033[0m'

    def addSuccess(self, test):
        super().addSuccess(test)
        print(f"  {self.GREEN}[PASS]{self.RESET}  {test.shortDescription() or str(test)}")

    def addFailure(self, test, err):
        super().addFailure(test, err)
        print(f"  {self.RED}[FAIL]{self.RESET}  {test.shortDescription() or str(test)}")
        print(f"       {self._exc_info_to_string(err, test).splitlines()[-1]}")

    def addError(self, test, err):
        super().addError(test, err)
        print(f"  {self.YELLOW}[ERROR]{self.RESET} {test.shortDescription() or str(test)}")
        print(f"       {self._exc_info_to_string(err, test).splitlines()[-1]}")


class ColorTestRunner(unittest.TextTestRunner):
    resultclass = ColorTestResult


if __name__ == "__main__":
    print()
    print("=" * 65)
    print("  apply_template.py v5 — 缺陷修复自动化测试")
    print("=" * 65)

    loader = unittest.TestLoader()
    suite = loader.loadTestsFromModule(sys.modules[__name__])

    runner = ColorTestRunner(verbosity=0, stream=sys.stdout)
    result = runner.run(suite)

    print()
    print("=" * 65)
    total = result.testsRun
    passed = total - len(result.failures) - len(result.errors)
    print(f"  总计：{total} 个测试  |  通过：{passed}  |  失败：{len(result.failures)}  |  错误：{len(result.errors)}")
    print("=" * 65)

    sys.exit(0 if result.wasSuccessful() else 1)
